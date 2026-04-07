# canvas_generator.py
# PM Tool - Business Model Canvas Generator
# Claude Architecture | OpenCode Implementation
# Path: C:\Users\USass\Billi_Bilder\PM_Tool\core\

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from datetime import date

COLORS = {
    "header_bg": RGBColor(0x1F, 0x49, 0x7D),
    "header_text": RGBColor(0xFF, 0xFF, 0xFF),
    "box_bg": RGBColor(0xF0, 0xF4, 0xF9),
    "box_border": RGBColor(0x1F, 0x49, 0x7D),
    "title_text": RGBColor(0x1F, 0x49, 0x7D),
    "body_text": RGBColor(0x1A, 0x1A, 0x1A),
    "footer_bg": RGBColor(0xDD, 0xE8, 0xF5),
}


def add_box(
    slide,
    left,
    top,
    width,
    height,
    title,
    content_lines,
    bg_color=None,
    title_color=None,
    body_color=None,
):
    bg = bg_color or COLORS["box_bg"]
    tc = title_color or COLORS["title_text"]
    bc = body_color or COLORS["body_text"]

    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = COLORS["box_border"]
    shape.line.width = Pt(1.2)

    tf = shape.text_frame
    tf.word_wrap = True

    p_title = tf.paragraphs[0]
    p_title.alignment = PP_ALIGN.LEFT
    run = p_title.add_run()
    run.text = title.upper()
    run.font.bold = True
    run.font.size = Pt(9)
    run.font.color.rgb = tc

    for line in content_lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = f"* {line}" if line.strip() else ""
        r.font.size = Pt(8)
        r.font.color.rgb = bc

    return shape


def create_canvas(data: dict, output_path: str = None):
    prs = Presentation()
    prs.slide_width = Cm(33.87)
    prs.slide_height = Cm(19.05)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    W = prs.slide_width
    H = prs.slide_height

    header = slide.shapes.add_shape(1, 0, 0, W, Cm(1.8))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS["header_bg"]
    header.line.fill.background()

    tf = header.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = tf.paragraphs[0].add_run()
    r.text = f"BUSINESS MODEL CANVAS - {data.get('project_name', 'Projekt').upper()}"
    r.font.bold = True
    r.font.size = Pt(14)
    r.font.color.rgb = COLORS["header_text"]

    margin = Cm(0.3)
    top_start = Cm(2.0)
    row_h = Cm(7.5)
    col_w = W / 5

    boxes_row1 = [
        ("Key Partners", data.get("key_partners", ["-"])),
        ("Key Activities", data.get("key_activities", ["-"])),
        ("Value Proposition", data.get("value_proposition", ["-"])),
        ("Customer Relations", data.get("customer_relations", ["-"])),
        ("Customer Segments", data.get("customer_segments", ["-"])),
    ]

    for i, (title, content) in enumerate(boxes_row1):
        add_box(
            slide,
            left=i * col_w + margin,
            top=top_start + margin,
            width=col_w - 2 * margin,
            height=row_h - 2 * margin,
            title=title,
            content_lines=content,
        )

    row2_top = top_start + row_h

    add_box(
        slide,
        left=col_w + margin,
        top=row2_top + margin,
        width=col_w - 2 * margin,
        height=row_h - 2 * margin,
        title="Key Resources",
        content_lines=data.get("key_resources", ["-"]),
    )

    add_box(
        slide,
        left=3 * col_w + margin,
        top=row2_top + margin,
        width=2 * col_w - 2 * margin,
        height=row_h - 2 * margin,
        title="Channels",
        content_lines=data.get("channels", ["-"]),
    )

    footer_top = top_start + 2 * row_h
    footer_h = H - footer_top - Cm(0.2)
    half_w = W / 2

    add_box(
        slide,
        left=margin,
        top=footer_top + margin,
        width=half_w - 2 * margin,
        height=footer_h - 2 * margin,
        title="Cost Structure",
        content_lines=data.get("cost_structure", ["-"]),
        bg_color=COLORS["footer_bg"],
    )

    add_box(
        slide,
        left=half_w + margin,
        top=footer_top + margin,
        width=half_w - 2 * margin,
        height=footer_h - 2 * margin,
        title="Revenue Streams",
        content_lines=data.get("revenue_streams", ["-"]),
        bg_color=COLORS["footer_bg"],
    )

    stamp = slide.shapes.add_textbox(W - Cm(5), H - Cm(0.7), Cm(4.8), Cm(0.6))
    tf2 = stamp.text_frame
    p = tf2.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r2 = p.add_run()
    r2.text = f"Erstellt: {date.today().strftime('%d.%m.%Y')} | {data.get('project_manager', '')}"
    r2.font.size = Pt(7)
    r2.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    if not output_path:
        safe = data.get("project_name", "Projekt").replace(" ", "_")
        output_path = f"C:\\Users\\USass\\Billi_Bilder\\PM_Tool\\{safe}_Canvas.pptx"

    prs.save(output_path)
    print(f"[OK] Canvas gespeichert: {output_path}")
    return output_path


if __name__ == "__main__":
    test_data = {
        "project_name": "PM Tool",
        "project_manager": "Udo Sassik",
        "key_partners": [
            "Ollama (lokale KI)",
            "OpenCode (Terminal-Agent)",
            "Claude (Strategie)",
            "GitHub (Versionskontrolle)",
        ],
        "key_activities": [
            "Charter Generator entwickeln",
            "Canvas Generator entwickeln",
            "Gantt Generator entwickeln",
            "KI-Assistent integrieren",
        ],
        "value_proposition": [
            "Automatische PM-Dokumente",
            "KI-gestuetzte Planung",
            "100% lokal & kostenlos",
            "Ein-Klick-Generierung",
        ],
        "customer_relations": [
            "Selbstnutzung (Udo)",
            "Zukuenftig: SGS/SPIE Teams",
            "GitHub Community",
        ],
        "customer_segments": [
            "Projektmanager",
            "Site Manager",
            "KMU ohne teure PM-Tools",
        ],
        "key_resources": [
            "Python-Stack",
            "Ollama llama3.1",
            "SHARED_BRAIN",
            "PM-Vorlagen",
        ],
        "channels": [
            "Lokal (CLI)",
            "Zukuenftig: Web-Interface",
            "GitHub Release",
        ],
        "cost_structure": [
            "0 Euro - nur kostenlose Tools",
            "Hardware (bestehendes Laptop)",
            "Zeit-Investition",
        ],
        "revenue_streams": [
            "Aktuell: Persoenlicher Nutzen",
            "Zukuenftig: Open Source Reputation",
            "Zukuenftig: Freelance PM-Services",
        ],
    }

    create_canvas(test_data)
