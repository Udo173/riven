#!/usr/bin/env python3
"""
Project Canvas PowerPoint Template Generator
Erstellt professionelle Präsentationen
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Farben
DARK_BLUE = RgbColor(10, 10, 30)
GOLD = RgbColor(255, 209, 102)
WHITE = RgbColor(255, 255, 255)
LIGHT_GRAY = RgbColor(200, 200, 200)


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ========== FOLIE 1: TITEL ==========
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Hintergrund
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # Titel
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "PROJECT CANVAS"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.CENTER

    # Untertitel
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Riven Musik-Projekt"
    p.font.size = Pt(32)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Datum
    date_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.5))
    tf = date_box.text_frame
    p = tf.paragraphs[0]
    p.text = "06.04.2026"
    p.font.size = Pt(18)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.CENTER

    # ========== FOLIE 2: ÜBERSICHT ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Hintergrund
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Project Canvas - Übersicht"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # 9 Boxen Grid
    elements = [
        ("GOALS", "Projektziele definieren"),
        ("USERS", "Zielgruppe identifizieren"),
        ("USER BENEFITS", "Nutzen für Anwender"),
        ("TEAM", "Projektteam zusammenstellen"),
        ("SCOPE", "Projektumfang abgrenzen"),
        ("STAKEHOLDERS", "Beteiligte erfassen"),
        ("DELIVERABLES", "Liefergegenstände"),
        ("ACTIVITIES", "Aktivitäten planen"),
        ("MILESTONES", "Meilensteine setzen"),
    ]

    positions = [
        (0.5, 1.3, 3.8, 1.5),
        (4.6, 1.3, 3.8, 1.5),
        (8.7, 1.3, 4, 1.5),
        (0.5, 3.1, 3.8, 1.5),
        (4.6, 3.1, 3.8, 1.5),
        (8.7, 3.1, 4, 1.5),
        (0.5, 4.9, 3.8, 1.5),
        (4.6, 4.9, 3.8, 1.5),
        (8.7, 4.9, 4, 1.5),
    ]

    for i, (title, desc) in enumerate(elements):
        x, y, w, h = positions[i]

        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(26, 26, 46)
        box.line.color.rgb = GOLD
        box.line.width = Pt(2)

        # Titel in Box
        title_tf = box.text_frame
        title_tf.paragraphs[0].text = title
        title_tf.paragraphs[0].font.size = Pt(18)
        title_tf.paragraphs[0].font.bold = True
        title_tf.paragraphs[0].font.color.rgb = GOLD
        title_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

        # Beschreibung
        p = title_tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    # ========== FOLIE 3: GOALS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "GOALS - Projektziele"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Inhalt
    goals = [
        "Professionelles Lyric-Video erstellen",
        "Songtext synchron zum Audio",
        "Motion Graphics mit Partikeln",
        "Cover Design mit Riven Logo",
        "AI-generierte Visualisierungen",
    ]

    content_box = slide.shapes.add_textbox(
        Inches(1), Inches(1.5), Inches(11), Inches(5)
    )
    tf = content_box.text_frame

    for i, goal in enumerate(goals):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✓  {goal}"
        p.font.size = Pt(28)
        p.font.color.rgb = WHITE
        p.space_after = Pt(20)

    # ========== FOLIE 4: USERS & BENEFITS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    # Titel
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "USERS & USER BENEFITS"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Zwei Spalten
    # Users
    users_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.5)
    )
    users_box.fill.solid()
    users_box.fill.fore_color.rgb = RgbColor(26, 26, 46)
    users_box.line.color.rgb = GOLD

    users_title = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(5), Inches(0.5)
    )
    tf = users_title.text_frame
    p = tf.paragraphs[0]
    p.text = "USERS (Zielgruppe)"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GOLD

    users = [
        "Riven Fans weltweit",
        "YouTube Zuschauer",
        "Spotify Hörer",
        "Musik-Enthusiasten",
        "AI & Tech Interessierte",
    ]
    users_content = slide.shapes.add_textbox(
        Inches(0.7), Inches(2.2), Inches(5), Inches(4)
    )
    tf = users_content.text_frame
    for i, user in enumerate(users):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {user}"
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(15)

    # Benefits
    benefits_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.3), Inches(6), Inches(5.5)
    )
    benefits_box.fill.solid()
    benefits_box.fill.fore_color.rgb = RgbColor(26, 26, 46)
    benefits_box.line.color.rgb = GOLD

    benefits_title = slide.shapes.add_textbox(
        Inches(6.7), Inches(1.5), Inches(5.5), Inches(0.5)
    )
    tf = benefits_title.text_frame
    p = tf.paragraphs[0]
    p.text = "USER BENEFITS"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GOLD

    benefits = [
        "Emotionaler Impact",
        "Inspiration zum Aufstehen",
        "Visuelle Schönheit",
        "Professionelle Qualität",
        "Kostenlos zugänglich",
    ]
    benefits_content = slide.shapes.add_textbox(
        Inches(6.7), Inches(2.2), Inches(5.5), Inches(4)
    )
    tf = benefits_content.text_frame
    for i, benefit in enumerate(benefits):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {benefit}"
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(15)

    # ========== FOLIE 5: TEAM & STAKEHOLDERS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "TEAM & STAKEHOLDERS"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOLD

    # Team Box
    team_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(5.5), Inches(5.5)
    )
    team_box.fill.solid()
    team_box.fill.fore_color.rgb = RgbColor(26, 26, 46)
    team_box.line.color.rgb = GOLD

    team_title = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.5), Inches(5), Inches(0.5)
    )
    tf = team_title.text_frame
    p = tf.paragraphs[0]
    p.text = "TEAM"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GOLD

    team_members = [
        "OpenCode (AI Agent)",
        "Claude (AI Agent)",
        "User (Udo Sassik)",
        "Billi (Arbeitsverzeichnis)",
    ]
    team_content = slide.shapes.add_textbox(
        Inches(0.7), Inches(2.2), Inches(5), Inches(4)
    )
    tf = team_content.text_frame
    for i, member in enumerate(team_members):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"👤 {member}"
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(15)

    # Stakeholders Box
    stake_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.5), Inches(1.3), Inches(6), Inches(5.5)
    )
    stake_box.fill.solid()
    stake_box.fill.fore_color.rgb = RgbColor(26, 26, 46)
    stake_box.line.color.rgb = GOLD

    stake_title = slide.shapes.add_textbox(
        Inches(6.7), Inches(1.5), Inches(5.5), Inches(0.5)
    )
    tf = stake_title.text_frame
    p = tf.paragraphs[0]
    p.text = "STAKEHOLDERS"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = GOLD

    stakeholders = [
        "Udo Sassik (Artist)",
        "Riven Community",
        "YouTube/Spotify Plattformen",
        "KI-Tool Entwickler",
    ]
    stake_content = slide.shapes.add_textbox(
        Inches(6.7), Inches(2.2), Inches(5.5), Inches(4)
    )
    tf = stake_content.text_frame
    for i, stake in enumerate(stakeholders):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {stake}"
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(15)

    # ========== FOLIE 6: DELIVERABLES ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "DELIVERABLES - Liefergegenstände"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOLD

    deliverables = [
        ("Motion Graphics Video", "Video mit Partikeln & Animation", "100%", "✅"),
        ("Cover Design", "Professionelles Cover mit Logo", "100%", "✅"),
        ("Lyric Video", "Songtext synchron zum Audio", "30%", "🔄"),
        ("AI Video", "KI-generierte Visualisierungen", "0%", "⏳"),
    ]

    y_pos = 1.5
    for name, desc, progress, status in deliverables:
        # Box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5),
            Inches(y_pos),
            Inches(12),
            Inches(1.2),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(26, 26, 46)
        box.line.color.rgb = GOLD

        # Name
        name_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(y_pos + 0.15), Inches(4), Inches(0.5)
        )
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{status} {name}"
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Description
        desc_box = slide.shapes.add_textbox(
            Inches(5), Inches(y_pos + 0.2), Inches(5), Inches(0.5)
        )
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_GRAY

        # Progress
        prog_box = slide.shapes.add_textbox(
            Inches(10.5), Inches(y_pos + 0.2), Inches(1.5), Inches(0.5)
        )
        tf = prog_box.text_frame
        p = tf.paragraphs[0]
        p.text = progress
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = GOLD
        p.alignment = PP_ALIGN.RIGHT

        y_pos += 1.4

    # ========== FOLIE 7: MILESTONES ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "MILESTONES - Meilensteine"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOLD

    milestones = [
        ("05.04.2026", "Motion Graphics", "✅"),
        ("06.04.2026", "Cover Design", "✅"),
        ("15.04.2026", "Lyric Video", "🔄"),
        ("01.05.2026", "AI Video", "⏳"),
        ("15.05.2026", "Final Compilation", "⏳"),
    ]

    # Timeline Linie
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2), Inches(3.5), Inches(9), Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = GOLD
    line.line.fill.background()

    x_pos = 1.5
    for date, name, status in milestones:
        # Punkt
        circle = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x_pos), Inches(3.25), Inches(0.5), Inches(0.5)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = GOLD if "✅" in status else RgbColor(100, 100, 100)
        circle.line.fill.background()

        # Datum drüber
        date_box = slide.shapes.add_textbox(
            Inches(x_pos - 0.3), Inches(2.5), Inches(1.5), Inches(0.5)
        )
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = date
        p.font.size = Pt(14)
        p.font.color.rgb = GOLD
        p.alignment = PP_ALIGN.CENTER

        # Name drunter
        name_box = slide.shapes.add_textbox(
            Inches(x_pos - 0.5), Inches(4), Inches(2), Inches(1)
        )
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"{status}\n{name}"
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        x_pos += 2.2

    # ========== FOLIE 8: RISKS ==========
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BLUE
    background.line.fill.background()

    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(12), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "RISKS - Risiken"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = GOLD

    risks = [
        ("Keine NVIDIA GPU", "KI-Video Rendering langsam", "🔴 Hoch"),
        ("Moviepy Windows", "Video-Codec Probleme", "🟡 Mittel"),
        ("Zeit-Limit", "Komplexität überschätzen", "🟡 Mittel"),
    ]

    y_pos = 1.5
    for name, desc, severity in risks:
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5),
            Inches(y_pos),
            Inches(12),
            Inches(1.5),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = RgbColor(26, 26, 46)
        box.line.color.rgb = RgbColor(255, 100, 100) if "🔴" in severity else GOLD

        name_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(y_pos + 0.2), Inches(8), Inches(0.5)
        )
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"⚠️ {name}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = WHITE

        desc_box = slide.shapes.add_textbox(
            Inches(0.7), Inches(y_pos + 0.8), Inches(8), Inches(0.5)
        )
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = LIGHT_GRAY

        sev_box = slide.shapes.add_textbox(
            Inches(10), Inches(y_pos + 0.4), Inches(2), Inches(0.5)
        )
        tf = sev_box.text_frame
        p = tf.paragraphs[0]
        p.text = severity
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RgbColor(255, 100, 100) if "🔴" in severity else GOLD
        p.alignment = PP_ALIGN.RIGHT

        y_pos += 1.8

    # Speichern
    output_path = (
        r"C:\Users\USass\Billi_Bilder\SHARED_BRAIN\Project_Canvas_Template.pptx"
    )
    prs.save(output_path)
    print(f"[OK] Praesentation erstellt: {output_path}")
    return output_path


if __name__ == "__main__":
    create_presentation()
